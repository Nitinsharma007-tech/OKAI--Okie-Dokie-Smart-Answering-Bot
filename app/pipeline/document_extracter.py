import json
import os
import re
import subprocess
import time
from pathlib import Path

import numpy as np
from dotenv import load_dotenv
from pptx import Presentation
from sentence_transformers import SentenceTransformer

from app.gemini_agent import GeminiAgent

load_dotenv()

# Optional: set this in your .env file if `pdftotext` isn't reliably on
# your system PATH, e.g.:
#   POPPLER_PATH=C:\poppler\poppler-26.02.0\Library\bin
# When unset, we just call "pdftotext" and rely on PATH as normal.
POPPLER_PATH = os.getenv("POPPLER_PATH", "")


class DocumentIngestor:
    """
    Turns raw manual/document files (.docx, .pdf, .txt) into structured
    knowledge topics and merges them into master_data/knowledge_master.json.

    Workflow:

        data/manual_uploads/*.docx / *.pdf / *.txt
                     │
                     ▼
            extract_text()           (pandoc / pdftotext)
                     │
                     ▼
            generate_topics()        (Gemini structures it into
                                       the same schema as the rest
                                       of knowledge_master.json)
                     │
                     ▼
            add_topics()             (dedupe by id, append)
                     │
                     ▼
            save()                   (writes knowledge_master.json)

    After running this, you still need to re-run your embedding
    pipeline (embedding_generator.py / knowledge_loader_embedding.py)
    and then main_cache_generator.py, since this script only updates
    knowledge_master.json itself.
    """

    SUPPORTED_EXTENSIONS = (".docx", ".pdf", ".pptx", ".txt")

    def __init__(
        self,
        input_folder="data/manual_uploads",
        master_file="master_data/knowledge_master.json",
        processed_folder="data/manual_uploads/processed",
    ):

        self.input_folder = Path(input_folder)
        self.master_file = Path(master_file)
        self.processed_folder = Path(processed_folder)

        self.input_folder.mkdir(parents=True, exist_ok=True)
        self.processed_folder.mkdir(parents=True, exist_ok=True)

        self.gemini = GeminiAgent()
        self.semantic_model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")

        self.master_data = {}
        self.topics = []

        self.load_master()

    # =====================================================
    # Load Existing Knowledge Master
    # =====================================================

    def load_master(self):

        if not self.master_file.exists():
            raise FileNotFoundError(f"\nKnowledge master not found:\n{self.master_file}")

        with open(self.master_file, "r", encoding="utf-8") as f:
            self.master_data = json.load(f)

        self.topics = self.master_data.get("topics", [])

        print(f"Loaded knowledge master with {len(self.topics)} existing topics.")

    # =====================================================
    # Extract Raw Text From A File
    # =====================================================

    def extract_text(self, file_path):

        file_path = Path(file_path)
        ext = file_path.suffix.lower()

        if ext == ".txt":
            return file_path.read_text(encoding="utf-8", errors="ignore")

        if ext == ".docx":

            result = subprocess.run(
                ["pandoc", "-t", "plain", str(file_path)],
                capture_output=True,
                encoding="utf-8",
                errors="replace",
                check=True,
            )
            return result.stdout

        if ext == ".pdf":

            pdftotext_exe = "pdftotext"

            if POPPLER_PATH:
                pdftotext_exe = str(Path(POPPLER_PATH) / "pdftotext.exe")

            result = subprocess.run(
                [pdftotext_exe, "-layout", "-enc", "UTF-8", str(file_path), "-"],
                capture_output=True,
                encoding="utf-8",
                errors="replace",
                check=True,
            )
            return result.stdout

        if ext == ".pptx":
            return self._extract_pptx_text(file_path)

        raise ValueError(f"Unsupported file type: {ext}")

    # =====================================================
    # Extract Text From A .pptx (Slide Text + Speaker Notes)
    # =====================================================

    def _extract_pptx_text(self, file_path):

        prs = Presentation(str(file_path))
        chunks = []

        for slide_num, slide in enumerate(prs.slides, start=1):

            chunks.append(f"--- Slide {slide_num} ---")

            for shape in slide.shapes:

                if shape.has_text_frame:

                    for paragraph in shape.text_frame.paragraphs:

                        text = "".join(run.text for run in paragraph.runs)

                        if text.strip():
                            chunks.append(text)

                if shape.has_table:

                    for row in shape.table.rows:

                        row_text = " | ".join(cell.text for cell in row.cells)

                        if row_text.strip():
                            chunks.append(row_text)

            # Speaker notes often contain the actual explanation behind
            # a slide's bullet points, so pull those in too.
            if slide.has_notes_slide:

                notes_text = slide.notes_slide.notes_text_frame.text

                if notes_text.strip():
                    chunks.append(f"[Speaker Notes] {notes_text}")

        return "\n".join(chunks)

    # =====================================================
    # Build The Gemini Prompt For Structuring
    # =====================================================

    def build_system_prompt(self):

        existing_modules = sorted({t.get("module", "") for t in self.topics if t.get("module")})

        return f"""
You are a knowledge-base structuring assistant for an ERP product's
support knowledge base.

You will be given the raw extracted text of a how-to document (it may
contain step-by-step instructions, screenshots captions, or general
explanations, possibly in a mix of English and Hindi).

Convert it into a JSON ARRAY of topic objects. Usually a document
covers ONE procedure, so return a single-item array — but if the
document clearly covers multiple distinct, unrelated procedures,
return one object per procedure.

Each object MUST use exactly this schema:

{{
  "module": "<short module/category name, reuse one of these existing modules if it fits: {', '.join(existing_modules)} — otherwise invent a short new one>",
  "topic": "<short topic title>",
  "summary": "<1-3 sentence summary of what this topic covers>",
  "navigation": ["<menu path items in order, if applicable>"],
  "steps": ["<ordered, clear, numbered-friendly steps>"],
  "business_rules": ["<any constraints/rules mentioned, or empty array>"],
  "important_notes": ["<any caveats/notes worth highlighting, or empty array>"],
  "questions": ["<4-6 natural questions a user might ask that this topic answers>"],
  "keywords": ["<5-8 relevant keywords>"],
  "related_topics": ["<related topic names if obvious, or empty array>"]
}}

Rules:
- Base the content ONLY on what's in the provided text. Do not invent steps or details that aren't there.
- Write steps as clear, standalone instructions (translate/paraphrase Hindi into English where present, don't duplicate both languages).
- Return ONLY the JSON array. No markdown fences, no preamble, no explanation.
"""

    # =====================================================
    # Generate Structured Topic(s) From Raw Text
    # =====================================================

    def generate_topics(self, raw_text, source_filename):

        system_prompt = self.build_system_prompt()
        user_prompt = f"DOCUMENT TEXT:\n\n{raw_text}"

        result = self.gemini.generate_json(system_prompt, user_prompt)

        # generate_json normally returns whatever Gemini outputs — make
        # sure we always end up with a list, even if it returned a
        # single object instead of an array.
        if isinstance(result, dict):
            result = [result]

        for topic in result:
            topic["id"] = self.make_id(topic["module"], topic["topic"])
            topic["source"] = {
                "file": source_filename,
                "generated_by": "Gemini 2.5 Flash (DocumentIngestor)",
                "version": "1.0",
            }
            topic["module_folder"] = topic["module"]
            topic["json_file"] = f"manual_{topic['id']}.json"

        return result

    # =====================================================
    # Slug / ID Helper
    # =====================================================

    def make_id(self, module, topic):

        slug = f"{module}_{topic}".lower()
        slug = re.sub(r"[^a-z0-9]+", "_", slug).strip("_")
        return slug

    # =====================================================
    # Add Topics (Dedupe By ID)
    # =====================================================

    def add_topics(self, new_topics):

        existing_ids = {t["id"] for t in self.topics}
        existing_questions = {
            " ".join(str(question).split()).lower()
            for topic in self.topics
            for question in topic.get("questions", [])
        }
        existing_topic_texts = []
        for topic in self.topics:
            chunks = [
                topic.get("module", ""),
                topic.get("topic", ""),
                topic.get("summary", ""),
                *topic.get("questions", []),
            ]
            combined = " ".join(str(part).strip() for part in chunks if str(part).strip())
            if combined:
                existing_topic_texts.append(combined)

        added = 0

        for topic in new_topics:

            topic_chunks = [
                topic.get("module", ""),
                topic.get("topic", ""),
                topic.get("summary", ""),
                *topic.get("questions", []),
            ]
            topic_text = " ".join(str(part).strip() for part in topic_chunks if str(part).strip())

            duplicate_by_semantics = False
            if topic_text:
                topic_embedding = self.semantic_model.encode(
                    topic_text,
                    convert_to_numpy=True,
                    normalize_embeddings=True,
                )
                for candidate_text in existing_topic_texts:
                    candidate_embedding = self.semantic_model.encode(
                        candidate_text,
                        convert_to_numpy=True,
                        normalize_embeddings=True,
                    )
                    score = float(np.dot(topic_embedding, candidate_embedding))
                    if score >= 0.92:
                        duplicate_by_semantics = True
                        break

            if topic["id"] in existing_ids:
                print(f"  ⚠ Skipping duplicate id: {topic['id']}")
                continue

            if duplicate_by_semantics:
                print(f"  ⚠ Skipping semantically duplicate topic: {topic.get('topic', 'Unknown')}")
                continue

            unique_questions = []
            for question in topic.get("questions", []):
                normalized = " ".join(str(question).split()).lower()
                if normalized and normalized not in existing_questions:
                    existing_questions.add(normalized)
                    unique_questions.append(question)
            topic["questions"] = unique_questions

            self.topics.append(topic)
            if topic_text:
                existing_topic_texts.append(topic_text)
            existing_ids.add(topic["id"])
            added += 1

        return added

    # =====================================================
    # Save Knowledge Master
    # =====================================================

    def save(self):

        module_counts = {}
        for t in self.topics:
            m = t.get("module", "General")
            module_counts[m] = module_counts.get(m, 0) + 1

        self.master_data["topics"] = self.topics
        self.master_data["modules"] = [
            {"module": m, "topics": c} for m, c in sorted(module_counts.items())
        ]
        self.master_data["total_modules"] = len(module_counts)
        self.master_data["total_topics"] = len(self.topics)
        self.master_data["generated_at"] = time.strftime("%Y-%m-%dT%H:%M:%S")

        with open(self.master_file, "w", encoding="utf-8") as f:
            json.dump(self.master_data, f, ensure_ascii=False, indent=2)

        print(f"✅ Saved knowledge master. Total topics: {len(self.topics)}")

    # =====================================================
    # Process A Single File
    # =====================================================

    def process_file(self, file_path):

        file_path = Path(file_path)

        print(f"\nProcessing: {file_path.name}")

        raw_text = self.extract_text(file_path)

        if not raw_text.strip():
            print("  ⚠ No extractable text found, skipping.")
            return 0

        topics = self.generate_topics(raw_text, file_path.name)
        added = self.add_topics(topics)

        print(f"  ✅ Added {added} topic(s) from this file.")

        # Move the processed file out of the inbox so re-runs don't
        # re-ingest the same document.
        destination = self.processed_folder / file_path.name
        file_path.rename(destination)

        return added

    # =====================================================
    # Process Every File In The Input Folder
    # =====================================================

    def run(self):

        files = [
            f for f in self.input_folder.iterdir()
            if f.is_file() and f.suffix.lower() in self.SUPPORTED_EXTENSIONS
        ]

        if not files:
            print(f"No files found in {self.input_folder}")
            return

        print("=" * 60)
        print(f"Found {len(files)} file(s) to ingest.")
        print("=" * 60)

        total_added = 0

        for f in files:

            try:
                total_added += self.process_file(f)

            except Exception as e:
                print(f"  ❌ Failed to process {f.name}: {e}")

            time.sleep(1)  # be gentle on the API

        self.save()

        print("=" * 60)
        print(f"🎉 Ingestion complete. {total_added} new topic(s) added across {len(files)} file(s).")
        print("=" * 60)